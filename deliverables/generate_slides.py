"""
Generate the Cotiviti Intern Assessment PowerPoint deck (.pptx).

Topic 1 - Clinical Natural Language Technology for Health Care.
POC: Chest X-Ray Report Generation System.

Author: Nitish Kumar - San Jose State University
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Palette ---------------------------------------------------------------
PURPLE = RGBColor(0x4B, 0x2E, 0x83)
PURPLE_D = RGBColor(0x32, 0x1E, 0x5A)
MAGENTA = RGBColor(0xC4, 0x00, 0x7A)
LIGHT = RGBColor(0xF4, 0xF1, 0xFA)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL = RGBColor(0x00, 0x9E, 0x8F)

FONT = "Calibri"
FONT_H = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def set_run(r, text, size, color, bold=False, italic=False, font=FONT):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font


def para(tf, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
         space_after=6, first=False, bullet=False, level=0, font=FONT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.level = level
    r = p.add_run()
    set_run(r, text, size, color, bold, italic, font)
    return p


def footer(s, page):
    tb, tf = textbox(s, Inches(0.4), Inches(7.05), Inches(9), Inches(0.35))
    para(tf, "Cotiviti Intern Assessment  ·  Nitish Kumar  ·  San Jose State University",
         9, GREY, first=True)
    tb2, tf2 = textbox(s, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.35))
    para(tf2, str(page), 9, GREY, align=PP_ALIGN.RIGHT, first=True)


def content_header(s, kicker, title):
    rect(s, 0, 0, SW, Inches(1.25), PURPLE)
    rect(s, 0, Inches(1.25), SW, Inches(0.06), MAGENTA)
    tb, tf = textbox(s, Inches(0.6), Inches(0.18), Inches(12), Inches(1.0))
    para(tf, kicker.upper(), 12, RGBColor(0xD9, 0xC7, 0xF0), bold=True, first=True, space_after=2)
    para(tf, title, 26, WHITE, bold=True, space_after=0)


def card(s, x, y, w, h, title, lines, accent=MAGENTA, title_color=PURPLE):
    c = rect(s, x, y, w, h, CARD)
    c.line.color.rgb = RGBColor(0xE2, 0xDD, 0xEE)
    c.line.width = Pt(1)
    bar = rect(s, x, y, Inches(0.09), h, accent)
    tb, tf = textbox(s, x + Inches(0.28), y + Inches(0.16), w - Inches(0.45), h - Inches(0.3))
    para(tf, title, 15, title_color, bold=True, first=True, space_after=6)
    for ln in lines:
        para(tf, ln, 11.5, DARK, space_after=4)
    return c


# ============================================================================
# SLIDE 1 — Title
# ============================================================================
s = slide()
rect(s, 0, 0, SW, SH, PURPLE)
rect(s, 0, 0, SW, SH, PURPLE)
# decorative bands
rect(s, 0, Inches(5.2), SW, Inches(0.09), MAGENTA)
rect(s, Inches(9.3), 0, Inches(4.05), SH, PURPLE_D)
rect(s, Inches(9.3), 0, Inches(0.09), SH, MAGENTA)

tb, tf = textbox(s, Inches(0.8), Inches(0.55), Inches(6), Inches(0.6))
para(tf, "COTIVITI", 20, WHITE, bold=True, first=True, space_after=0)
para(tf, "INTERN ASSESSMENT", 12, RGBColor(0xD9, 0xC7, 0xF0), bold=True, space_after=0)

tb, tf = textbox(s, Inches(0.8), Inches(2.1), Inches(8.2), Inches(2.6))
para(tf, "Clinical Natural Language", 40, WHITE, bold=True, first=True, space_after=0)
para(tf, "Technology for Health Care", 40, WHITE, bold=True, space_after=10)
para(tf, "Grounding NLP, Computer Vision & Multimodal AI in Radiology",
     18, RGBColor(0xE8, 0xDE, 0xF6), italic=True, space_after=0)

tb, tf = textbox(s, Inches(0.8), Inches(5.5), Inches(8), Inches(1.4))
para(tf, "Nitish Kumar", 20, WHITE, bold=True, first=True, space_after=2)
para(tf, "San Jose State University", 14, RGBColor(0xD9, 0xC7, 0xF0), space_after=2)
para(tf, "Topic 1  ·  Report + Proof of Concept", 12, RGBColor(0xB9, 0xA6, 0xD8), space_after=0)

# right column highlights
tb, tf = textbox(s, Inches(9.7), Inches(2.2), Inches(3.3), Inches(4))
para(tf, "POC", 13, MAGENTA, bold=True, first=True, space_after=2)
para(tf, "Chest X-Ray Report", 15, WHITE, bold=True, space_after=0)
para(tf, "Generation System", 15, WHITE, bold=True, space_after=14)
para(tf, "DenseNet-121 + Transformer", 11.5, RGBColor(0xE8, 0xDE, 0xF6), space_after=4)
para(tf, "RAG Clinical Agent", 11.5, RGBColor(0xE8, 0xDE, 0xF6), space_after=4)
para(tf, "FastAPI · Gradio · Docker", 11.5, RGBColor(0xE8, 0xDE, 0xF6), space_after=4)

# ============================================================================
# SLIDE 2 — The Problem
# ============================================================================
s = slide()
content_header(s, "The Opportunity", "80% of Health Data Is Unstructured")
footer(s, 2)
tb, tf = textbox(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(1.1))
para(tf, "The details that matter (diagnoses, procedures, and the proof behind every claim) sit in "
         "free-text notes, scanned faxes, and images. Computers cannot search, code, or audit them directly.",
     15, GREY, first=True)

data = [
    ("Text", "Physician notes, discharge summaries, pathology reports", MAGENTA),
    ("Scanned", "Faxed orders and referrals digitized with OCR", PURPLE),
    ("Images", "X-ray, CT and MRI read by computer vision", TEAL),
    ("Codes", "The clean ICD / CPT codes payment integrity needs", RGBColor(0xE0,0x7A,0x00)),
]
x = Inches(0.6)
cw = Inches(2.95)
gap = Inches(0.18)
for title, txt, col in data:
    card(s, x, Inches(3.0), cw, Inches(2.6), title, [txt], accent=col, title_color=col)
    x = Emu(int(x) + int(cw) + int(gap))

tb, tf = textbox(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.9))
para(tf, "Clinical NLP turns all four into structured data. That is the foundation for accurate coding, "
         "payment integrity, and quality measurement.", 13, PURPLE, bold=True, first=True)

# ============================================================================
# SLIDE 3 — Trends / Evolution
# ============================================================================
s = slide()
content_header(s, "Trends", "Past → Present → Future")
footer(s, 3)
eras = [
    ("PAST", "2000-2017", "Rules and Word Lists",
     ["cTAKES, MetaMap map text to UMLS / SNOMED / ICD",
      "Hand-written rules and medical word lists",
      "Accurate, but fragile and costly to update"], PURPLE),
    ("PRESENT", "2018-2023", "Learning from Big Data",
     ["BioBERT, ClinicalBERT, GatorTron models",
      "Image models near expert level: DenseNet, ViT",
      "Good at finding terms and hiding patient IDs"], MAGENTA),
    ("FUTURE", "2024 on", "Generative and Multimodal",
     ["Clinical LLMs (Med-PaLM 2) and LMMs",
      "One model reads an X-ray and drafts the report",
      "RAG keeps answers grounded, cuts made-up facts"], TEAL),
]
x = Inches(0.6)
cw = Inches(3.95)
gap = Inches(0.2)
for kicker, yrs, title, lines, col in eras:
    top = rect(s, x, Inches(1.7), cw, Inches(0.95), col)
    tb, tf = textbox(s, x + Inches(0.25), Inches(1.78), cw - Inches(0.4), Inches(0.85),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, kicker + "   " + yrs, 12, WHITE, bold=True, first=True, space_after=1)
    para(tf, title, 15, WHITE, bold=True, space_after=0)
    card(s, x, Inches(2.75), cw, Inches(3.4), "", lines, accent=col, title_color=col)
    x = Emu(int(x) + int(cw) + int(gap))

# arrows
for ax in [Inches(4.62), Inches(8.77)]:
    a = s.shapes.add_shape(MSO_SHAPE.CHEVRON, ax, Inches(2.0), Inches(0.35), Inches(0.55))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0xC9,0xBD,0xE0); a.line.fill.background()
    a.shadow.inherit = False

# ============================================================================
# SLIDE 4 — Opportunities & Threats
# ============================================================================
s = slide()
content_header(s, "Analysis", "Opportunities & Threats")
footer(s, 4)
card(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(4.9), "OPPORTUNITIES",
     ["Automated coding: find miscoding and missing documentation at scale",
      "Report drafting from images: faster radiology and pathology turnaround",
      "Traceable AI: RAG shows the exact source behind every answer",
      "More reviewer time: automation handles routine cases so experts focus on the hard ones"],
     accent=TEAL, title_color=TEAL)
card(s, Inches(6.9), Inches(1.7), Inches(6.0), Inches(4.9), "THREATS AND RISKS",
     ["Made-up facts: models can state false things that sound convincing",
      "Privacy: patient data needs HIPAA-grade handling and strict controls",
      "Bias and drift: uneven results, and accuracy drops as coding rules change",
      "Regulation: FDA and payer oversight of clinical AI keeps changing"],
     accent=MAGENTA, title_color=MAGENTA)

# ============================================================================
# SLIDE 5 — Recommendations for Cotiviti
# ============================================================================
s = slide()
content_header(s, "Strategy", "Recommended Actions for Cotiviti")
footer(s, 5)
recs = [
    ("01", "Grounded Clinical-NLP Layer",
     "Combine medical language models (ClinicalBERT / GatorTron) with a trusted RAG library of coding "
     "policies and guidelines, so every AI code links back to an official source."),
    ("02", "Multimodal Document Understanding",
     "Use OCR, NLP, and computer vision together to check images and notes against submitted claims, "
     "catching gaps that text-only review misses."),
    ("03", "People in the Loop",
     "Track model versions, watch for drift, and send low-confidence cases to a person, so automation "
     "supports expert reviewers and stays compliant."),
]
y = Inches(1.7)
for num, title, txt in recs:
    badge = rect(s, Inches(0.6), y, Inches(1.0), Inches(1.35), PURPLE, shape=MSO_SHAPE.OVAL)
    tb, tf = textbox(s, Inches(0.6), y, Inches(1.0), Inches(1.35), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, num, 24, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    c = card(s, Inches(1.85), y, Inches(11.0), Inches(1.35), title, [txt], accent=MAGENTA)
    y = Emu(int(y) + int(Inches(1.55)))

# ============================================================================
# SLIDE 6 — POC Overview
# ============================================================================
s = slide()
content_header(s, "Proof of Concept", "Chest X-Ray Report Generation System")
footer(s, 6)
tb, tf = textbox(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.8))
para(tf, "A full working demo: from a raw X-ray to a written report and a grounded clinical Q&A agent.",
     15, GREY, first=True)

# pipeline boxes
steps = [
    ("X-Ray Input", "JPEG · PNG · DICOM", PURPLE),
    ("DenseNet-121", "Vision encoder → 256-d", MAGENTA),
    ("Transformer", "Decoder → report text", PURPLE),
    ("RAG Agent", "ClinicalBERT + ChromaDB", TEAL),
    ("Clinical Q&A", "ReAct reasoning trace", RGBColor(0xE0,0x7A,0x00)),
]
x = Inches(0.6)
bw = Inches(2.28)
gap = Inches(0.14)
for title, sub, col in steps:
    b = rect(s, x, Inches(2.6), bw, Inches(1.5), col)
    tb, tf = textbox(s, x + Inches(0.12), Inches(2.7), bw - Inches(0.24), Inches(1.3),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, 14, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=4)
    para(tf, sub, 10.5, RGBColor(0xF0,0xEA,0xF8), align=PP_ALIGN.CENTER, space_after=0)
    x = Emu(int(x) + int(bw) + int(gap))

card(s, Inches(0.6), Inches(4.5), Inches(6.0), Inches(2.15), "What It Shows",
     ["Computer vision plus NLP writes the report",
      "RAG keeps the AI's answers grounded and traceable",
      "DICOM windowing (lung / mediastinum / bone)"],
     accent=PURPLE, title_color=PURPLE)
card(s, Inches(6.9), Inches(4.5), Inches(6.0), Inches(2.15), "Production-Style Stack",
     ["FastAPI + Celery + Redis async inference",
      "Prometheus + OpenTelemetry observability",
      "Docker + GitHub Actions · Gradio 4-tab UI"],
     accent=TEAL, title_color=TEAL)

# ============================================================================
# SLIDE 7 — System Architecture (HLD)
# ============================================================================
import os
_HLD = os.path.join(os.path.dirname(__file__), "hld_architecture.png")
s = slide()
content_header(s, "How It Works", "System Architecture (High-Level Design)")
footer(s, 7)
if os.path.exists(_HLD):
    # size by height so it fits between header and footer, then center horizontally
    pic_h = Inches(4.95)
    pic_w = Emu(int(pic_h * 13.2 / 7.1))
    left = Emu(int((SW - pic_w) / 2))
    s.shapes.add_picture(_HLD, left, Inches(1.5), height=pic_h)
tb, tf = textbox(s, Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.35))
para(tf, "Four layers: clients to FastAPI to services (model, DICOM, RAG, async) to data, "
         "with observability and delivery across the whole stack.",
     11, GREY, align=PP_ALIGN.CENTER, first=True)

# ============================================================================
# SLIDE 8 — API at a Glance
# ============================================================================
s = slide()
content_header(s, "Interface", "The API at a Glance")
footer(s, 8)
tb, tf = textbox(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.6))
para(tf, "A REST API built with FastAPI. Full interactive docs are auto-generated at /docs (Swagger).",
     14, GREY, first=True)
api_groups = [
    ("SYSTEM", PURPLE, [
        ("GET  /health", "liveness and model-loaded status"),
        ("GET  /metadata", "model architecture details"),
        ("GET  /metrics", "Prometheus metrics"),
    ]),
    ("INFERENCE", MAGENTA, [
        ("POST /predict", "report from JPEG / PNG"),
        ("POST /predict/dicom", "report from a DICOM file"),
        ("POST /preview/dicom", "windowed DICOM preview image"),
        ("GET  /task/{id}", "poll an async job"),
    ]),
    ("AGENT", TEAL, [
        ("POST /agent/query", "grounded question answering"),
        ("POST /agent/search", "semantic search over reports"),
        ("POST /agent/index", "add reports to the vector store"),
        ("GET  /agent/stats", "vector store statistics"),
    ]),
]
x = Inches(0.6)
cw = Inches(4.05)
gap = Inches(0.13)
for name, col, rows in api_groups:
    hdr = rect(s, x, Inches(2.25), cw, Inches(0.55), col)
    tb, tf = textbox(s, x + Inches(0.2), Inches(2.28), cw - Inches(0.3), Inches(0.5),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, name, 13, WHITE, bold=True, first=True)
    c = rect(s, x, Inches(2.8), cw, Inches(3.7), CARD)
    c.line.color.rgb = RGBColor(0xE2, 0xDD, 0xEE)
    c.line.width = Pt(1)
    tb, tf = textbox(s, x + Inches(0.22), Inches(2.95), cw - Inches(0.4), Inches(3.5))
    first = True
    for ep, desc in rows:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(9)
        r = p.add_run(); set_run(r, ep, 11.5, col, bold=True, font="Consolas")
        r2 = p.add_run(); set_run(r2, "\n" + desc, 10.5, DARK)
    x = Emu(int(x) + int(cw) + int(gap))

# ============================================================================
# SLIDE 9 — Demo / How it maps
# ============================================================================
s = slide()
content_header(s, "Live Demo", "What You'll See in the Screenshare")
footer(s, 9)
demo = [
    ("Report Generation", "Upload an X-ray and get a generated radiology report", MAGENTA),
    ("Clinical Agent", "Ask a clinical question, see the step-by-step reasoning and a grounded answer", PURPLE),
    ("Knowledge Base", "Add reports, search by meaning, and manage the vector store", TEAL),
    ("System and Telemetry", "Live health, Prometheus metrics, OpenTelemetry status", RGBColor(0xE0,0x7A,0x00)),
]
y = Inches(1.8)
for title, txt, col in demo:
    dot = rect(s, Inches(0.7), Emu(int(y)+int(Inches(0.12))), Inches(0.22), Inches(0.22), col, shape=MSO_SHAPE.OVAL)
    tb, tf = textbox(s, Inches(1.15), y, Inches(11.6), Inches(1.1))
    para(tf, title, 16, col, bold=True, first=True, space_after=2)
    para(tf, txt, 13, DARK, space_after=0)
    y = Emu(int(y) + int(Inches(1.18)))

# ============================================================================
# SLIDE 8 — Closing
# ============================================================================
s = slide()
rect(s, 0, 0, SW, SH, PURPLE)
rect(s, 0, Inches(5.2), SW, Inches(0.09), MAGENTA)
tb, tf = textbox(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(3),)
para(tf, "Key Takeaways", 34, WHITE, bold=True, first=True, space_after=18)
para(tf, "Clinical NLP unlocks the 80% of health data trapped in text and images.",
     18, RGBColor(0xE8,0xDE,0xF6), space_after=12)
para(tf, "Grounded, multimodal AI with people in the loop is the safe path for payment integrity.",
     18, RGBColor(0xE8,0xDE,0xF6), space_after=12)
para(tf, "The POC proves the full pipeline works, from model to deployable service.",
     18, RGBColor(0xE8,0xDE,0xF6), space_after=12)

tb, tf = textbox(s, Inches(0.9), Inches(5.5), Inches(11), Inches(1.2))
para(tf, "Thank You", 24, WHITE, bold=True, first=True, space_after=2)
para(tf, "Nitish Kumar  ·  San Jose State University  ·  nitish.kumar@sjsu.edu",
     13, RGBColor(0xD9,0xC7,0xF0), space_after=0)

import os
out = os.path.join(os.path.dirname(__file__), "Nitish_Kumar_Presentation.pptx")
prs.save(out)
print("Saved:", out, "-", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
